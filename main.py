import streamlit as st
from PyPDF2 import PdfMerger
from pptx import Presentation
from pathlib import Path
from io import BytesIO
import subprocess
from streamlit_sortables import sort_items

# PPT 파일을 PPTX로 변환하는 함수
def convert_ppt_to_pptx(input_file_path, output_dir):
    output_file = Path(output_dir) / f"{Path(input_file_path).stem}.pptx"
    subprocess.run([
        "libreoffice", 
        "--headless", 
        "--convert-to", "pptx", 
        "--outdir", str(output_dir), 
        str(input_file_path)
    ])
    return output_file

# App Title
st.title("📎 PPT & PDF 결합 도구 (By 석리송)")

# File Upload
uploaded_files = st.file_uploader(
    "📤 PPT, PPTX, PDF 파일을 업로드하세요", 
    type=["ppt", "pptx", "pdf"], 
    accept_multiple_files=True
)

if uploaded_files:
    # 파일 이름 목록
    filenames = [file.name for file in uploaded_files]

    # 기본 파일 순서 (오름차순 정렬)
    filenames = sorted(filenames)

    # Drag-and-Drop으로 파일 순서 변경
    st.write("### 파일 순서를 Drag-and-Drop으로 변경하세요:")
    sorted_filenames = sort_items(filenames)

    st.write("#### 선택된 파일 순서:")
    st.write(sorted_filenames)

    # 출력 파일 이름 설정
    ppt_output_name = st.text_input(
        "📁 PPT 결합 파일 이름 (기본값: merged.pptx)", 
        value="merged.pptx"
    )
    pdf_output_name = st.text_input(
        "📁 PDF 결합 파일 이름 (기본값: merged.pdf)", 
        value="merged.pdf"
    )

    # 결합 버튼
    if st.button("결합하기"):
        temp_dir = Path("temp_files")
        temp_dir.mkdir(exist_ok=True)

        # PPT 결합 처리
        try:
            merged_presentation = Presentation()
            for filename in sorted_filenames:
                file = next(f for f in uploaded_files if f.name == filename)
                temp_file_path = temp_dir / file.name

                # 파일 임시 저장
                with open(temp_file_path, "wb") as temp_file:
                    temp_file.write(file.read())

                # PPT 파일을 PPTX로 변환
                if temp_file_path.suffix == ".ppt":
                    temp_file_path = convert_ppt_to_pptx(temp_file_path, temp_dir)

                # PPTX 파일 읽기 및 슬라이드 결합
                presentation = Presentation(temp_file_path)
                for slide in presentation.slides:
                    blank_slide_layout = merged_presentation.slide_layouts[6]
                    slide_copy = merged_presentation.slides.add_slide(blank_slide_layout)
                    for shape in slide.shapes:
                        el = shape.element
                        slide_copy.shapes._spTree.insert_element_before(el, 'p:extLst')

            # 저장
            ppt_output_path = temp_dir / ppt_output_name
            merged_presentation.save(ppt_output_path)
            with open(ppt_output_path, "rb") as f:
                st.download_button(
                    "📥 PPT 결합 파일 다운로드",
                    f,
                    file_name=ppt_output_name,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
        except Exception as e:
            st.error(f"PPT 결합 중 오류 발생: {e}")

        # PDF 결합 처리
        try:
            merger = PdfMerger()
            for filename in sorted_filenames:
                file = next(f for f in uploaded_files if f.name == filename)
                if file.type == "application/pdf":
                    merger.append(BytesIO(file.read()))
            pdf_output_path = temp_dir / pdf_output_name
            merger.write(pdf_output_path)
            merger.close()
            with open(pdf_output_path, "rb") as f:
                st.download_button(
                    "📥 PDF 결합 파일 다운로드",
                    f,
                    file_name=pdf_output_name,
                    mime="application/pdf"
                )
        except Exception as e:
            st.error(f"PDF 결합 중 오류 발생: {e}")

        # 임시 파일 정리
        for temp_file in temp_dir.iterdir():
            temp_file.unlink()
        temp_dir.rmdir()
else:
    st.warning("최소 하나의 파일을 업로드해주세요.")
